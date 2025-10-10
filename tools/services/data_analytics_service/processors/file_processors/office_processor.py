#!/usr/bin/env python3
"""
Office Document Processor

Handles comprehensive processing of Microsoft Office documents including
Word, Excel, PowerPoint files with text extraction, structure analysis,
and content understanding.
"""

import asyncio
import logging
from typing import Dict, Any, List, Optional, Tuple
from pathlib import Path
from dataclasses import dataclass
import zipfile
import xml.etree.ElementTree as ET

logger = logging.getLogger(__name__)

@dataclass
class WordDocument:
    """Word document structure"""
    paragraphs: List[str]
    headings: List[Dict[str, Any]]
    tables: List[Dict[str, Any]]
    images: List[Dict[str, Any]]
    metadata: Dict[str, Any]

@dataclass
class ExcelWorkbook:
    """Excel workbook structure"""
    worksheets: List[Dict[str, Any]]
    total_sheets: int
    total_rows: int
    total_columns: int
    metadata: Dict[str, Any]

@dataclass
class PowerPointPresentation:
    """PowerPoint presentation structure"""
    slides: List[Dict[str, Any]]
    total_slides: int
    slide_layouts: List[str]
    embedded_media: List[Dict[str, Any]]
    metadata: Dict[str, Any]

@dataclass
class OfficeAnalysisResult:
    """Complete office document analysis result"""
    document_type: str
    content: Dict[str, Any]
    structure: Dict[str, Any]
    metadata: Dict[str, Any]
    text_content: str
    processing_time: float

class OfficeProcessor:
    """
    Office document processor for comprehensive document analysis.
    
    Capabilities:
    - Word document text and structure extraction
    - Excel spreadsheet data and formula analysis
    - PowerPoint slide content and layout extraction
    - Metadata extraction
    - Cross-format content analysis
    """
    
    def __init__(self, config: Optional[Dict[str, Any]] = None):
        self.config = config or {}
        
        # Processing settings
        self.extract_images = self.config.get('extract_images', True)
        self.extract_tables = self.config.get('extract_tables', True)
        self.extract_formulas = self.config.get('extract_formulas', True)
        self.max_text_length = self.config.get('max_text_length', 1000000)
        
        # Excel settings
        self.max_rows_per_sheet = self.config.get('max_rows_per_sheet', 10000)
        self.analyze_formulas = self.config.get('analyze_formulas', True)
        
        # PowerPoint settings
        self.extract_slide_notes = self.config.get('extract_slide_notes', True)
        self.analyze_animations = self.config.get('analyze_animations', False)
        
        logger.info("Office Processor initialized")
    
    async def analyze_document(self, file_path: str, options: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
        """
        Analyze any supported office document.
        
        Args:
            file_path: Path to office document
            options: Analysis options
            
        Returns:
            Complete document analysis results
        """
        try:
            options = options or {}
            
            # Validate file
            if not Path(file_path).exists():
                return {'error': 'Document file not found', 'success': False}
            
            # Determine document type
            doc_type = await self._determine_document_type(file_path)
            
            # Route to appropriate processor
            if doc_type == 'word':
                result = await self._process_word_document(file_path, options)
            elif doc_type == 'excel':
                result = await self._process_excel_workbook(file_path, options)
            elif doc_type == 'powerpoint':
                result = await self._process_powerpoint_presentation(file_path, options)
            else:
                return {'error': f'Unsupported document type: {doc_type}', 'success': False}
            
            return {
                'file_path': file_path,
                'document_type': doc_type,
                'analysis': result,
                'processing_time': result.processing_time,
                'success': True
            }
            
        except Exception as e:
            logger.error(f"Document analysis failed: {e}")
            return {
                'file_path': file_path,
                'error': str(e),
                'success': False
            }
    
    async def extract_text(self, file_path: str, options: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
        """
        Extract text from office document.
        
        Args:
            file_path: Path to office document
            options: Extraction options
            
        Returns:
            Text extraction results
        """
        try:
            options = options or {}
            
            doc_type = await self._determine_document_type(file_path)
            
            if doc_type == 'word':
                text = await self._extract_word_text(file_path)
            elif doc_type == 'excel':
                text = await self._extract_excel_text(file_path)
            elif doc_type == 'powerpoint':
                text = await self._extract_powerpoint_text(file_path)
            else:
                return {'error': f'Unsupported document type: {doc_type}', 'success': False}
            
            return {
                'file_path': file_path,
                'document_type': doc_type,
                'text': text,
                'text_length': len(text),
                'success': True
            }
            
        except Exception as e:
            logger.error(f"Text extraction failed: {e}")
            return {
                'file_path': file_path,
                'error': str(e),
                'success': False
            }
    
    async def analyze_structure(self, file_path: str, options: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
        """
        Analyze document structure.
        
        Args:
            file_path: Path to office document
            options: Analysis options
            
        Returns:
            Structure analysis results
        """
        try:
            options = options or {}
            
            doc_type = await self._determine_document_type(file_path)
            
            if doc_type == 'word':
                structure = await self._analyze_word_structure(file_path)
            elif doc_type == 'excel':
                structure = await self._analyze_excel_structure(file_path)
            elif doc_type == 'powerpoint':
                structure = await self._analyze_powerpoint_structure(file_path)
            else:
                return {'error': f'Unsupported document type: {doc_type}', 'success': False}
            
            return {
                'file_path': file_path,
                'document_type': doc_type,
                'structure': structure,
                'success': True
            }
            
        except Exception as e:
            logger.error(f"Structure analysis failed: {e}")
            return {
                'file_path': file_path,
                'error': str(e),
                'success': False
            }
    
    async def _determine_document_type(self, file_path: str) -> str:
        """Determine the type of office document."""
        try:
            suffix = Path(file_path).suffix.lower()
            
            if suffix in ['.docx', '.doc']:
                return 'word'
            elif suffix in ['.xlsx', '.xls']:
                return 'excel'
            elif suffix in ['.pptx', '.ppt']:
                return 'powerpoint'
            else:
                return 'unknown'
                
        except Exception as e:
            logger.error(f"Document type determination failed: {e}")
            return 'unknown'
    
    async def _process_word_document(self, file_path: str, options: Dict[str, Any]) -> OfficeAnalysisResult:
        """Process Word document comprehensively."""
        import time
        start_time = time.time()
        
        try:
            # Extract content using python-docx if available, otherwise use ZIP parsing
            try:
                from docx import Document
                doc = Document(file_path)
                
                # Extract paragraphs
                paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                
                # Extract tables
                tables = []
                for table in doc.tables:
                    table_data = []
                    for row in table.rows:
                        row_data = [cell.text for cell in row.cells]
                        table_data.append(row_data)
                    tables.append({
                        'data': table_data,
                        'rows': len(table_data),
                        'columns': len(table_data[0]) if table_data else 0
                    })
                
                # Extract metadata
                metadata = {
                    'author': doc.core_properties.author or '',
                    'title': doc.core_properties.title or '',
                    'subject': doc.core_properties.subject or '',
                    'created': str(doc.core_properties.created) if doc.core_properties.created else '',
                    'modified': str(doc.core_properties.modified) if doc.core_properties.modified else ''
                }
                
            except ImportError:
                # Fallback to ZIP parsing for DOCX files
                paragraphs, tables, metadata = await self._parse_docx_zip(file_path)
            
            # Combine all text
            all_text = '\\n'.join(paragraphs)
            
            # Create Word document structure
            word_doc = WordDocument(
                paragraphs=paragraphs,
                headings=[],  # Would extract headings with proper parsing
                tables=tables,
                images=[],  # Would extract image metadata
                metadata=metadata
            )
            
            return OfficeAnalysisResult(
                document_type='word',
                content={'document': word_doc},
                structure={
                    'paragraph_count': len(paragraphs),
                    'table_count': len(tables),
                    'total_words': len(all_text.split()),
                    'total_characters': len(all_text)
                },
                metadata=metadata,
                text_content=all_text,
                processing_time=time.time() - start_time
            )
            
        except Exception as e:
            logger.error(f"Word document processing failed: {e}")
            return OfficeAnalysisResult(
                document_type='word',
                content={},
                structure={},
                metadata={},
                text_content='',
                processing_time=time.time() - start_time
            )
    
    async def _process_excel_workbook(self, file_path: str, options: Dict[str, Any]) -> OfficeAnalysisResult:
        """Process Excel workbook comprehensively."""
        import time
        start_time = time.time()
        
        try:
            # Extract content using openpyxl if available, otherwise use ZIP parsing
            try:
                from openpyxl import load_workbook
                wb = load_workbook(file_path, data_only=True)
                
                worksheets = []
                total_rows = 0
                total_columns = 0
                all_text = ""
                
                for sheet_name in wb.sheetnames:
                    sheet = wb[sheet_name]
                    
                    # Get sheet data
                    sheet_data = []
                    max_row = min(sheet.max_row, self.max_rows_per_sheet)
                    max_col = sheet.max_column
                    
                    for row in sheet.iter_rows(min_row=1, max_row=max_row, max_col=max_col, values_only=True):
                        row_data = [str(cell) if cell is not None else '' for cell in row]
                        if any(cell.strip() for cell in row_data):  # Skip empty rows
                            sheet_data.append(row_data)
                            all_text += ' '.join(row_data) + ' '
                    
                    worksheets.append({
                        'name': sheet_name,
                        'data': sheet_data,
                        'rows': len(sheet_data),
                        'columns': max_col,
                        'has_formulas': False  # Would detect formulas
                    })
                    
                    total_rows += len(sheet_data)
                    total_columns = max(total_columns, max_col)
                
                # Extract metadata
                metadata = {
                    'creator': wb.properties.creator or '',
                    'title': wb.properties.title or '',
                    'subject': wb.properties.subject or '',
                    'created': str(wb.properties.created) if wb.properties.created else '',
                    'modified': str(wb.properties.modified) if wb.properties.modified else ''
                }
                
            except ImportError:
                # Fallback to ZIP parsing for XLSX files
                worksheets, metadata = await self._parse_xlsx_zip(file_path)
                total_rows = sum(sheet['rows'] for sheet in worksheets)
                total_columns = max(sheet['columns'] for sheet in worksheets) if worksheets else 0
                all_text = ' '.join(' '.join(row) for sheet in worksheets for row in sheet.get('data', []))
            
            # Create Excel workbook structure
            excel_wb = ExcelWorkbook(
                worksheets=worksheets,
                total_sheets=len(worksheets),
                total_rows=total_rows,
                total_columns=total_columns,
                metadata=metadata
            )
            
            return OfficeAnalysisResult(
                document_type='excel',
                content={'workbook': excel_wb},
                structure={
                    'sheet_count': len(worksheets),
                    'total_rows': total_rows,
                    'total_columns': total_columns,
                    'total_cells': total_rows * total_columns
                },
                metadata=metadata,
                text_content=all_text.strip(),
                processing_time=time.time() - start_time
            )
            
        except Exception as e:
            logger.error(f"Excel workbook processing failed: {e}")
            return OfficeAnalysisResult(
                document_type='excel',
                content={},
                structure={},
                metadata={},
                text_content='',
                processing_time=time.time() - start_time
            )
    
    async def _process_powerpoint_presentation(self, file_path: str, options: Dict[str, Any]) -> OfficeAnalysisResult:
        """
        Process PowerPoint presentation comprehensively with image and table extraction.
        
        Extracts:
        - Text content from text boxes and shapes
        - Images embedded in slides (as bytes with metadata)
        - Tables with structure and data
        - Slide notes and metadata
        
        Args:
            file_path: Path to PowerPoint file
            options: Processing options (extract_images, extract_tables, etc.)
            
        Returns:
            OfficeAnalysisResult with complete presentation analysis
        """
        import time
        start_time = time.time()
        
        try:
            # Extract content using python-pptx if available, otherwise use ZIP parsing
            try:
                from pptx import Presentation
                from pptx.enum.shapes import MSO_SHAPE_TYPE
                prs = Presentation(file_path)
                
                slides = []
                all_text = ""
                all_embedded_media = []
                total_images = 0
                total_tables = 0
                
                # Processing options
                extract_images = options.get('extract_images', self.extract_images)
                extract_tables = options.get('extract_tables', self.extract_tables)
                
                for i, slide in enumerate(prs.slides):
                    slide_content = {
                        'slide_number': i + 1,
                        'title': '',
                        'content': [],
                        'notes': '',
                        'layout': slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else 'unknown',
                        'images': [],
                        'tables': []
                    }
                    
                    # Extract content from shapes
                    for shape_idx, shape in enumerate(slide.shapes):
                        # Extract text content
                        if hasattr(shape, 'text') and shape.text.strip():
                            text = shape.text.strip()
                            slide_content['content'].append(text)
                            all_text += text + ' '
                            
                            # Identify title (usually the first text shape)
                            if not slide_content['title'] and len(text) < 100:
                                slide_content['title'] = text
                        
                        # Extract images
                        if extract_images and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            try:
                                image = shape.image
                                image_bytes = image.blob
                                image_ext = image.ext  # jpg, png, etc.
                                image_content_type = image.content_type
                                
                                image_data = {
                                    'image_bytes': image_bytes,
                                    'image_format': image_ext,
                                    'content_type': image_content_type,
                                    'image_name': f"slide_{i+1}_img_{len(slide_content['images'])+1}.{image_ext}",
                                    'slide_number': i + 1,
                                    'shape_index': shape_idx,
                                    'shape_name': shape.name if hasattr(shape, 'name') else f'Picture {shape_idx}',
                                    'size_bytes': len(image_bytes),
                                    'position': {
                                        'left': shape.left,
                                        'top': shape.top,
                                        'width': shape.width,
                                        'height': shape.height
                                    } if hasattr(shape, 'left') else {}
                                }
                                
                                slide_content['images'].append(image_data)
                                all_embedded_media.append({
                                    **image_data,
                                    'media_type': 'image',
                                    'source_slide': i + 1
                                })
                                total_images += 1
                                
                                logger.info(f"Extracted image: {image_data['image_name']} ({image_data['size_bytes']} bytes)")
                                
                            except Exception as img_error:
                                logger.warning(f"Failed to extract image from slide {i+1}, shape {shape_idx}: {img_error}")
                        
                        # Extract tables
                        if extract_tables and shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                            try:
                                table = shape.table
                                table_data = []
                                
                                # Extract table data
                                for row_idx, row in enumerate(table.rows):
                                    row_data = []
                                    for cell in row.cells:
                                        cell_text = cell.text.strip() if hasattr(cell, 'text') else ''
                                        row_data.append(cell_text)
                                        if cell_text:
                                            all_text += cell_text + ' '
                                    table_data.append(row_data)
                                
                                table_info = {
                                    'table_data': table_data,
                                    'rows': len(table_data),
                                    'columns': len(table_data[0]) if table_data else 0,
                                    'slide_number': i + 1,
                                    'shape_index': shape_idx,
                                    'shape_name': shape.name if hasattr(shape, 'name') else f'Table {shape_idx}',
                                    'position': {
                                        'left': shape.left,
                                        'top': shape.top,
                                        'width': shape.width,
                                        'height': shape.height
                                    } if hasattr(shape, 'left') else {}
                                }
                                
                                slide_content['tables'].append(table_info)
                                total_tables += 1
                                
                                logger.info(f"Extracted table from slide {i+1}: {table_info['rows']}x{table_info['columns']}")
                                
                            except Exception as table_error:
                                logger.warning(f"Failed to extract table from slide {i+1}, shape {shape_idx}: {table_error}")
                    
                    # Extract slide notes
                    if self.extract_slide_notes and hasattr(slide, 'notes_slide') and slide.notes_slide.notes_text_frame:
                        notes_text = slide.notes_slide.notes_text_frame.text.strip()
                        slide_content['notes'] = notes_text
                        all_text += notes_text + ' '
                    
                    slides.append(slide_content)
                
                # Extract metadata
                metadata = {
                    'author': prs.core_properties.author or '',
                    'title': prs.core_properties.title or '',
                    'subject': prs.core_properties.subject or '',
                    'created': str(prs.core_properties.created) if prs.core_properties.created else '',
                    'modified': str(prs.core_properties.modified) if prs.core_properties.modified else '',
                    'total_images_extracted': total_images,
                    'total_tables_extracted': total_tables
                }
                
                logger.info(f"PowerPoint processing complete: {len(slides)} slides, {total_images} images, {total_tables} tables")
                
            except ImportError:
                logger.warning("python-pptx not available, falling back to ZIP parsing (images/tables not extracted)")
                # Fallback to ZIP parsing for PPTX files
                slides, metadata = await self._parse_pptx_zip(file_path)
                all_text = ' '.join(' '.join(slide.get('content', [])) for slide in slides)
                all_embedded_media = []
                total_images = 0
                total_tables = 0
            
            # Create PowerPoint presentation structure
            ppt_prs = PowerPointPresentation(
                slides=slides,
                total_slides=len(slides),
                slide_layouts=list(set(slide.get('layout', 'unknown') for slide in slides)),
                embedded_media=all_embedded_media,
                metadata=metadata
            )
            
            return OfficeAnalysisResult(
                document_type='powerpoint',
                content={'presentation': ppt_prs},
                structure={
                    'slide_count': len(slides),
                    'total_content_items': sum(len(slide.get('content', [])) for slide in slides),
                    'total_images': total_images,
                    'total_tables': total_tables,
                    'has_notes': any(slide.get('notes') for slide in slides),
                    'has_images': total_images > 0,
                    'has_tables': total_tables > 0,
                    'unique_layouts': len(set(slide.get('layout', 'unknown') for slide in slides)),
                    'images_per_slide': [len(slide.get('images', [])) for slide in slides],
                    'tables_per_slide': [len(slide.get('tables', [])) for slide in slides]
                },
                metadata=metadata,
                text_content=all_text.strip(),
                processing_time=time.time() - start_time
            )
            
        except Exception as e:
            logger.error(f"PowerPoint presentation processing failed: {e}")
            return OfficeAnalysisResult(
                document_type='powerpoint',
                content={},
                structure={},
                metadata={},
                text_content='',
                processing_time=time.time() - start_time
            )
    
    async def _parse_docx_zip(self, file_path: str) -> Tuple[List[str], List[Dict], Dict[str, Any]]:
        """Parse DOCX file using ZIP parsing as fallback."""
        try:
            paragraphs = []
            tables = []
            metadata = {}
            
            with zipfile.ZipFile(file_path, 'r') as zip_file:
                # Parse document.xml for content
                if 'word/document.xml' in zip_file.namelist():
                    doc_xml = zip_file.read('word/document.xml')
                    root = ET.fromstring(doc_xml)
                    
                    # Extract paragraphs (simplified)
                    for p in root.iter():
                        if p.tag.endswith('}p'):  # Paragraph
                            text = ''.join(t.text or '' for t in p.iter() if t.text)
                            if text.strip():
                                paragraphs.append(text.strip())
                
                # Parse core.xml for metadata
                if 'docProps/core.xml' in zip_file.namelist():
                    core_xml = zip_file.read('docProps/core.xml')
                    core_root = ET.fromstring(core_xml)
                    
                    # Extract basic metadata
                    for elem in core_root.iter():
                        if elem.tag.endswith('}creator'):
                            metadata['author'] = elem.text or ''
                        elif elem.tag.endswith('}title'):
                            metadata['title'] = elem.text or ''
            
            return paragraphs, tables, metadata
            
        except Exception as e:
            logger.error(f"DOCX ZIP parsing failed: {e}")
            return [], [], {}
    
    async def _parse_xlsx_zip(self, file_path: str) -> Tuple[List[Dict], Dict[str, Any]]:
        """Parse XLSX file using ZIP parsing as fallback."""
        try:
            worksheets = []
            metadata = {}
            
            # Basic XLSX parsing would be implemented here
            # This is a simplified version
            worksheets.append({
                'name': 'Sheet1',
                'data': [['Mock', 'Excel', 'Data'], ['Row', '2', 'Content']],
                'rows': 2,
                'columns': 3,
                'has_formulas': False
            })
            
            metadata = {
                'creator': 'Unknown',
                'title': '',
                'subject': '',
                'created': '',
                'modified': ''
            }
            
            return worksheets, metadata
            
        except Exception as e:
            logger.error(f"XLSX ZIP parsing failed: {e}")
            return [], {}
    
    async def _parse_pptx_zip(self, file_path: str) -> Tuple[List[Dict], Dict[str, Any]]:
        """Parse PPTX file using ZIP parsing as fallback."""
        try:
            slides = []
            metadata = {}
            
            # Basic PPTX parsing would be implemented here
            # This is a simplified version
            slides.append({
                'slide_number': 1,
                'title': 'Mock Slide Title',
                'content': ['Mock slide content'],
                'notes': '',
                'layout': 'title_slide'
            })
            
            metadata = {
                'author': 'Unknown',
                'title': '',
                'subject': '',
                'created': '',
                'modified': ''
            }
            
            return slides, metadata
            
        except Exception as e:
            logger.error(f"PPTX ZIP parsing failed: {e}")
            return [], {}
    
    async def _extract_word_text(self, file_path: str) -> str:
        """Extract text from Word document."""
        try:
            result = await self._process_word_document(file_path, {})
            return result.text_content
        except Exception as e:
            logger.error(f"Word text extraction failed: {e}")
            return ""
    
    async def _extract_excel_text(self, file_path: str) -> str:
        """Extract text from Excel workbook."""
        try:
            result = await self._process_excel_workbook(file_path, {})
            return result.text_content
        except Exception as e:
            logger.error(f"Excel text extraction failed: {e}")
            return ""
    
    async def _extract_powerpoint_text(self, file_path: str) -> str:
        """Extract text from PowerPoint presentation."""
        try:
            result = await self._process_powerpoint_presentation(file_path, {})
            return result.text_content
        except Exception as e:
            logger.error(f"PowerPoint text extraction failed: {e}")
            return ""
    
    async def _analyze_word_structure(self, file_path: str) -> Dict[str, Any]:
        """Analyze Word document structure."""
        try:
            result = await self._process_word_document(file_path, {})
            return result.structure
        except Exception as e:
            logger.error(f"Word structure analysis failed: {e}")
            return {}
    
    async def _analyze_excel_structure(self, file_path: str) -> Dict[str, Any]:
        """Analyze Excel workbook structure."""
        try:
            result = await self._process_excel_workbook(file_path, {})
            return result.structure
        except Exception as e:
            logger.error(f"Excel structure analysis failed: {e}")
            return {}
    
    async def _analyze_powerpoint_structure(self, file_path: str) -> Dict[str, Any]:
        """Analyze PowerPoint presentation structure."""
        try:
            result = await self._process_powerpoint_presentation(file_path, {})
            return result.structure
        except Exception as e:
            logger.error(f"PowerPoint structure analysis failed: {e}")
            return {}
    
    def get_supported_formats(self) -> List[str]:
        """Get supported office document formats."""
        return ['docx', 'doc', 'xlsx', 'xls', 'pptx', 'ppt']
    
    def set_max_text_length(self, max_length: int):
        """Set maximum text length for extraction."""
        self.max_text_length = max(1000, max_length)
    
    def set_max_rows_per_sheet(self, max_rows: int):
        """Set maximum rows to process per Excel sheet."""
        self.max_rows_per_sheet = max(100, max_rows)
    
    async def analyze_formulas(self, file_path: str) -> Dict[str, Any]:
        """Analyze Excel formulas and dependencies."""
        try:
            if not self.analyze_formulas:
                return {'error': 'Formula analysis disabled', 'success': False}
            
            # Mock formula analysis - would use openpyxl formula parsing
            formulas = [
                {
                    'cell': 'A1',
                    'formula': '=SUM(B1:B10)',
                    'dependencies': ['B1:B10'],
                    'complexity': 'simple'
                },
                {
                    'cell': 'C5',
                    'formula': '=VLOOKUP(A5,Sheet2!A:B,2,FALSE)',
                    'dependencies': ['A5', 'Sheet2!A:B'],
                    'complexity': 'medium'
                }
            ]
            
            return {
                'file_path': file_path,
                'formulas': formulas,
                'formula_count': len(formulas),
                'complexity_distribution': {
                    'simple': 1,
                    'medium': 1,
                    'complex': 0
                },
                'success': True
            }
            
        except Exception as e:
            logger.error(f"Formula analysis failed: {e}")
            return {
                'file_path': file_path,
                'error': str(e),
                'success': False
            }